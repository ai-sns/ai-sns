# -*- coding: utf-8 -*-
"""
Document Loader - Extract text from various file formats
"""
import logging
from pathlib import Path
from typing import Optional
import pypdf
from docx import Document
from pptx import Presentation
import openpyxl
import os
from runtime.shared import debug_info

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Load and extract text from various document formats"""

    @staticmethod
    def _convert_legacy_office(file_path: Path) -> Optional[Path]:
        suffix = file_path.suffix.lower()
        if suffix not in {'.doc', '.ppt', '.xls'}:
            return None

        converted = file_path.with_name(f"{file_path.stem}.__converted__{suffix}x")

        try:
            import win32com.client  # type: ignore

            if suffix == '.doc':
                word = win32com.client.DispatchEx('Word.Application')
                word.Visible = False
                doc = None
                try:
                    doc = word.Documents.Open(str(file_path))
                    doc.SaveAs(str(converted), FileFormat=16)
                finally:
                    if doc is not None:
                        doc.Close(False)
                    word.Quit()

            elif suffix == '.ppt':
                powerpoint = win32com.client.DispatchEx('PowerPoint.Application')
                presentation = None
                try:
                    presentation = powerpoint.Presentations.Open(str(file_path), WithWindow=False)
                    presentation.SaveAs(str(converted), FileFormat=24)
                finally:
                    if presentation is not None:
                        presentation.Close()
                    powerpoint.Quit()

            elif suffix == '.xls':
                excel = win32com.client.DispatchEx('Excel.Application')
                excel.Visible = False
                excel.DisplayAlerts = False
                wb = None
                try:
                    wb = excel.Workbooks.Open(str(file_path))
                    wb.SaveAs(str(converted), FileFormat=51)
                finally:
                    if wb is not None:
                        wb.Close(False)
                    excel.Quit()

            if converted.exists() and converted.stat().st_size > 0:
                return converted

            return None
        except Exception as e:
            logger.error(f"Legacy Office conversion failed for {file_path}: {e}")
            return None

    @staticmethod
    def load_pdf(file_path: Path) -> str:
        """Extract text from PDF file"""
        try:
            text = []
            with open(file_path, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
            return '\n\n'.join(text)
        except Exception as e:
            logger.error(f"Error loading PDF {file_path}: {e}")
            raise

    @staticmethod
    def load_docx(file_path: Path) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            return '\n\n'.join(text)
        except Exception as e:
            logger.error(f"Error loading DOCX {file_path}: {e}")
            raise

    @staticmethod
    def load_pptx(file_path: Path) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return '\n\n'.join(text)
        except Exception as e:
            logger.error(f"Error loading PPTX {file_path}: {e}")
            raise

    @staticmethod
    def load_xlsx(file_path: Path) -> str:
        """Extract text from XLSX file"""
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text = []
            for sheet in wb.worksheets:
                text.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text.append(row_text)
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"Error loading XLSX {file_path}: {e}")
            raise

    @staticmethod
    def load_txt(file_path: Path) -> str:
        """Load text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Error loading TXT {file_path}: {e}")
                raise
        except Exception as e:
            logger.error(f"Error loading TXT {file_path}: {e}")
            raise

    @staticmethod
    def load_md(file_path: Path) -> str:
        return DocumentLoader.load_txt(file_path)

    @classmethod
    def load_document(cls, file_path: Path) -> Optional[str]:
        """Load document based on file extension"""
        suffix = file_path.suffix.lower()

        converted_path: Optional[Path] = None
        if suffix in {'.doc', '.ppt', '.xls'}:
            converted_path = cls._convert_legacy_office(file_path)
            if converted_path:
                file_path = converted_path
                suffix = file_path.suffix.lower()

        loaders = {
            '.pdf': cls.load_pdf,
            '.docx': cls.load_docx,
            '.pptx': cls.load_pptx,
            '.xlsx': cls.load_xlsx,
            '.txt': cls.load_txt,
            '.md': cls.load_md,
            '.markdown': cls.load_md,
        }

        loader = loaders.get(suffix)
        if loader:
            try:
                return loader(file_path)
            except Exception as e:
                logger.error(f"Error loading document {file_path}: {e}")
                return None
            finally:
                if converted_path:
                    try:
                        os.remove(converted_path)
                    except Exception:
                        pass
        else:
            logger.warning(f"Unsupported file format: {suffix}")
            if converted_path:
                try:
                    os.remove(converted_path)
                except Exception:
                    pass
            return None
